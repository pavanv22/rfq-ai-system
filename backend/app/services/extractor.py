"""
Extractor service - extracts text from various document formats.

Supported formats:
- PDF (.pdf)
- Word Documents (.docx, .doc)
- PowerPoint Presentations (.pptx, .ppt)
- Excel Spreadsheets (.xlsx, .xls)
- Plain Text (.txt)
- CSV Files (.csv)
- Images (.png, .jpg, .jpeg, .bmp - uses OCR)
"""

import pdfplumber
from docx import Document
import pytesseract
from PIL import Image
import os
import json
from typing import Optional
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    from openpyxl import load_workbook
except ImportError:
    openpyxl = None
    load_workbook = None


def get_file_extension(file_path: str) -> str:
    """Get the file extension in lowercase."""
    return Path(file_path).suffix.lower()


def extract_text(file_path: str) -> Optional[str]:
    """
    Extract text from various document formats.
    
    Args:
        file_path: Path to the document file
    
    Returns:
        Extracted text or None if extraction fails
    """
    if not os.path.exists(file_path):
        return None
    
    ext = get_file_extension(file_path)
    
    try:
        if ext == ".pdf":
            return extract_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return extract_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            return extract_pptx(file_path)
        elif ext in [".xlsx", ".xls"]:
            return extract_excel(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".gif"]:
            return extract_image(file_path)
        elif ext == ".txt":
            return extract_text_file(file_path)
        elif ext == ".csv":
            return extract_csv_file(file_path)
        else:
            return None  # Unsupported format
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return None


def extract_pdf(file_path: str) -> str:
    """Extract text from PDF files."""
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text += f"\n[Page {i+1}]\n"
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"Error extracting PDF: {e}")
    
    return text


def extract_docx(file_path: str) -> str:
    """Extract text from Word documents."""
    text = ""
    try:
        doc = Document(file_path)
        
        # Extract from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text += " | ".join(row_text) + "\n"
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
    
    return text


def extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentations."""
    text = ""
    
    if Presentation is None:
        print("python-pptx not installed. Cannot extract PPTX files.")
        return text
    
    try:
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"\n[Slide {slide_num}]\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text += " | ".join(row_text) + "\n"
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
    
    return text


def extract_excel(file_path: str) -> str:
    """Extract text from Excel spreadsheets."""
    text = ""
    
    if load_workbook is None:
        print("openpyxl not installed. Cannot extract Excel files.")
        return text
    
    try:
        workbook = load_workbook(file_path)
        
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            text += f"\n[Sheet: {sheet_name}]\n"
            
            for row in worksheet.iter_rows(values_only=True):
                row_text = []
                for cell_value in row:
                    if cell_value is not None:
                        row_text.append(str(cell_value).strip())
                if row_text:
                    text += " | ".join(row_text) + "\n"
    except Exception as e:
        print(f"Error extracting Excel: {e}")
    
    return text


def extract_image(file_path: str) -> str:
    """Extract text from images using OCR (Tesseract)."""
    text = ""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
    except Exception as e:
        print(f"Error extracting image (ensure Tesseract is installed): {e}")
    
    return text


def extract_text_file(file_path: str) -> str:
    """Extract text from plain text files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try different encoding if UTF-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file: {e}")
            return ""
    except Exception as e:
        print(f"Error extracting text file: {e}")
        return ""


def extract_csv_file(file_path: str) -> str:
    """Extract content from CSV files as formatted text."""
    try:
        import csv
        text = ""
        with open(file_path, 'r', encoding='utf-8') as f:
            csv_reader = csv.reader(f)
            for row_idx, row in enumerate(csv_reader):
                text += " | ".join(row) + "\n"
        return text
    except Exception as e:
        print(f"Error extracting CSV file: {e}")
        return ""


def extract_with_metadata(file_path: str) -> dict:
    """
    Extract text and metadata from a document.
    
    Returns:
        Dictionary with extracted text and metadata
    """
    if not os.path.exists(file_path):
        return {"error": "File not found", "file_path": file_path}
    
    file_name = os.path.basename(file_path)
    file_ext = get_file_extension(file_path)
    file_size = os.path.getsize(file_path)
    
    extracted_text = extract_text(file_path)
    
    return {
        "file_name": file_name,
        "file_extension": file_ext,
        "file_size_bytes": file_size,
        "extracted_text": extracted_text,
        "text_length": len(extracted_text) if extracted_text else 0,
        "extraction_status": "success" if extracted_text else "failed"
    }
